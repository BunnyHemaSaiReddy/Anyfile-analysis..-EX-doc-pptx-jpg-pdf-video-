from PyPDF2 import PdfReader
import docx
from pptx import Presentation
import pandas as pd
from moviepy.editor import VideoFileClip
import PIL.Image as Image
import speech_recognition as sr
import assemblyai as aai
import io

def transcribe_audio(audio_file_path):
    try:
        aai.settings.api_key = "28108f8b185a472c8062dc14b789b18e"
        transcriber = aai.Transcriber()
        transcript = transcriber.transcribe(audio_file_path)
        #st.write(transcript.text)
        return transcript.text
    except Exception:
        pass

def process_file(uploaded_file):
    file_type = uploaded_file.type
    
    if file_type.startswith('image/'):
        img = Image.open(uploaded_file)
        return img, 'image'
    
    elif file_type == 'application/pdf':
        pdf_reader = PdfReader(uploaded_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        return text, 'pdf'
    
    elif file_type in ['application/msword', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document']:
        doc = docx.Document(uploaded_file)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text, 'doc'
    
    elif file_type in ['application/vnd.ms-powerpoint', 'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
        ppt = Presentation(uploaded_file)
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text, 'ppt'
    
    elif file_type in ['application/vnd.ms-excel', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
        df = pd.read_excel(uploaded_file)
        text = df.to_string()
        return text, 'excel'
    
    # elif file_type.startswith('video/'):
    #     #st.video(uploaded_file)
    #     video = VideoFileClip(io.BytesIO(uploaded_file.read()))
    #     audio = video.audio
    #     audio_buffer = io.BytesIO()
    #     audio.write_audiofile(audio_buffer)
    #     audio_buffer.seek(0)
    #     transcript = transcribe_audio(audio_buffer)
    #     return transcript, 'video'
    
    elif file_type.startswith('audio/'):
        #st.audio(uploaded_file)
        audio_buffer = io.BytesIO(uploaded_file.read())
        transcript = transcribe_audio(audio_buffer)
        return transcript, 'audio'
